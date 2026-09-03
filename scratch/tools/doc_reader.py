#!/usr/bin/env python3
"""
doc_reader.py - PowerPoint (.pptx) & PDF (.pdf) Document Inspection & Extraction Tool
Supports text extraction, slide notes, table parsing, and PDF-to-image conversion for Antigravity & Jetson.

Usage:
  python3 doc_reader.py <file.pptx | file.pdf>
  python3 doc_reader.py <file.pptx | file.pdf> --output <output.md>
  python3 doc_reader.py <file.pdf> --to-images [--dpi 150]
"""

import sys
import os
import argparse
import subprocess
from typing import Optional

def inspect_pptx(file_path: str, output_path: Optional[str] = None) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    lines = []
    lines.append(f"# 📊 PowerPoint Presentation Inspection: {os.path.basename(file_path)}")
    lines.append(f"• Total Slides: {len(prs.slides)}\n")

    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Slide {idx}")
        
        # Extract Title
        title = slide.shapes.title.text.strip() if slide.shapes.title else "Untitled Slide"
        lines.append(f"### Title: {title}\n")

        # Extract Shapes & Text
        body_texts = []
        tables = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                txt = shape.text.strip()
                if txt:
                    body_texts.append(txt)
            if shape.has_table:
                table_lines = []
                for row in shape.table.rows:
                    row_txt = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                    table_lines.append("| " + " | ".join(row_txt) + " |")
                if table_lines:
                    # Markdown table header separator
                    sep = "| " + " | ".join(["---"] * len(shape.table.columns)) + " |"
                    table_lines.insert(1, sep)
                    tables.append("\n".join(table_lines))

        if body_texts:
            lines.append("**Content:**")
            for b in body_texts:
                for sub_line in b.split('\n'):
                    sub_line = sub_line.strip()
                    if sub_line:
                        lines.append(f"- {sub_line}")
            lines.append("")

        if tables:
            lines.append("**Tables:**")
            for tbl in tables:
                lines.append(tbl)
                lines.append("")

        # Extract Notes Slide
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"**Speaker Notes:**\n> {notes}\n")

        lines.append("---\n")

    result = "\n".join(lines)
    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(result)
        print(f"✅ Extracted PowerPoint markdown saved to: {output_path}")
    return result


def inspect_pdf(file_path: str, output_path: Optional[str] = None, to_images: bool = False, dpi: int = 150) -> str:
    import pypdf

    reader = pypdf.PdfReader(file_path)
    total_pages = len(reader.pages)
    lines = []
    lines.append(f"# 📄 PDF Document Inspection: {os.path.basename(file_path)}")
    lines.append(f"• Total Pages: {total_pages}\n")

    meta = reader.metadata
    if meta:
        lines.append("## Metadata")
        if meta.title:
            lines.append(f"- Title: {meta.title}")
        if meta.author:
            lines.append(f"- Author: {meta.author}")
        if meta.subject:
            lines.append(f"- Subject: {meta.subject}")
        lines.append("")

    for idx, page in enumerate(reader.pages, start=1):
        lines.append(f"## Page {idx} / {total_pages}")
        text = page.extract_text()
        if text and text.strip():
            lines.append(text.strip())
        else:
            lines.append("*(No text extracted - scanned page or pure image)*")
        lines.append("\n---\n")

    result = "\n".join(lines)
    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(result)
        print(f"✅ Extracted PDF text markdown saved to: {output_path}")

    # Optional: Convert pages to image files using pdftoppm
    if to_images:
        out_prefix = os.path.splitext(file_path)[0] + "_pages"
        cmd = ["pdftoppm", "-png", "-r", str(dpi), file_path, out_prefix]
        try:
            subprocess.run(cmd, check=True)
            print(f"✅ Rendered PDF pages as PNG images with prefix: {out_prefix}")
        except Exception as e:
            print(f"⚠️ Failed to render PDF images via pdftoppm: {e}")

    return result


def main():
    parser = argparse.ArgumentParser(description="Inspect & extract PowerPoint (.pptx) and PDF (.pdf) documents.")
    parser.add_argument("file", help="Path to .pptx or .pdf document")
    parser.add_argument("--output", "-o", help="Optional output markdown file path")
    parser.add_argument("--to-images", action="store_true", help="Render PDF pages as PNG images via pdftoppm")
    parser.add_argument("--dpi", type=int, default=150, help="DPI for image rendering (default: 150)")
    args = parser.parse_args()

    if not os.path.exists(args.file):
        print(f"❌ Error: File not found: {args.file}", file=sys.stderr)
        sys.exit(1)

    ext = os.path.splitext(args.file)[1].lower()
    if ext == ".pptx":
        content = inspect_pptx(args.file, args.output)
        if not args.output:
            print(content)
    elif ext == ".pdf":
        content = inspect_pdf(args.file, args.output, args.to_images, args.dpi)
        if not args.output:
            print(content)
    else:
        print(f"❌ Error: Unsupported file format '{ext}'. Only .pptx and .pdf are supported.", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
